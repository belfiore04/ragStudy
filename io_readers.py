import io
import os
import shutil
from pathlib import Path
from typing import List, Optional
from PIL import Image
from langchain.schema import Document
from config import PDF_RENDER_DPI



def read_pdf(file_bytes: bytes, name: str) -> List[Document]:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(file_bytes))
    docs: List[Document] = []
    for i, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if text:
            docs.append(Document(page_content=text, metadata={"source": name, "page": i}))
    return docs

def read_pptx(file_bytes: bytes, name: str) -> List[Document]:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    docs: List[Document] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts: List[str] = []
        for shp in slide.shapes:
            if hasattr(shp, "shapes"):
                for s2 in getattr(shp, "shapes", []):
                    if getattr(s2, "has_text_frame", False) and s2.text_frame:
                        txt = "\n".join(p.text for p in s2.text_frame.paragraphs).strip();
                        if txt: parts.append(txt)
            if getattr(shp, "has_text_frame", False) and shp.text_frame:
                txt = "\n".join(p.text for p in shp.text_frame.paragraphs).strip();
                if txt: parts.append(txt)
            if getattr(shp, "has_table", False) and shp.table:
                cells = []
                for r in shp.table.rows:
                    for c in r.cells:
                        t = (c.text or "").strip()
                        if t: cells.append(t)
                if cells: parts.append("\n".join(cells))
        text = "\n".join(t for t in parts if t).strip()
        if text:
            docs.append(Document(page_content=text, metadata={"source": name, "slide": i}))
    return docs
def read_docx(file_bytes: bytes, name: str) -> List[Document]:
    import docx
    doc = docx.Document(io.BytesIO(file_bytes))
    paras = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
    text = "\n".join(paras)
    return [Document(page_content=text, metadata={"source": name})] if text else []




def read_txt(file_bytes: bytes, name: str) -> List[Document]:
    text = file_bytes.decode("utf-8", errors="ignore").strip()
    return [Document(page_content=text, metadata={"source": name})] if text else []

def pdf_page_to_image(pdf_path: Path, page_number: int) -> Optional[Image.Image]:
    try:
        import fitz # PyMuPDF
        doc = fitz.open(pdf_path)
        if 1 <= page_number <= len(doc):
            page = doc.load_page(page_number - 1)
            pix = page.get_pixmap(dpi=PDF_RENDER_DPI)
            return Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
    except Exception:
        pass
    try:
        from pdf2image import convert_from_path
        imgs = convert_from_path(str(pdf_path), dpi=PDF_RENDER_DPI, first_page=page_number, last_page=page_number)
        return imgs[0] if imgs else None
    except Exception:
        return None




def convert_to_pdf_with_libreoffice(input_path: Path, out_dir: Path) -> Optional[Path]:
    if shutil.which("soffice") is None:
        return None
    out_dir.mkdir(parents=True, exist_ok=True)
    cmd = f'"{shutil.which("soffice")}" --headless --convert-to pdf --outdir "{out_dir}" "{input_path}"'
    code = os.system(cmd)
    if code != 0:
        return None
    out_pdf = out_dir / (input_path.stem + ".pdf")
    return out_pdf if out_pdf.exists() else None